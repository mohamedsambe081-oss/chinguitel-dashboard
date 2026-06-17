from __future__ import annotations

from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile

# IMPORTANT: Django runs in a web/server context, not in a desktop GUI loop.
# Force Matplotlib to use a non-GUI backend to avoid Tkinter thread errors on Windows
# such as: RuntimeError: main thread is not in main loop / Tcl_AsyncDelete.
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import PercentFormatter
from django.conf import settings
from django.core.files import File
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from django.utils import timezone

from revenue.models import DataUpload, PowerPointReport
from .analytics import (
    category_breakdown,
    calendar_heatmap,
    detect_anomalies,
    filter_records,
    forecast_4_weeks,
    package_breakdown,
    summary,
    time_series,
    eda,
    detail_dashboard,
)

BLUE = "0070C0"
VIOLET = "7030A0"
WHITE = "FFFFFF"
LIGHT = "F8FBFF"
DARK = "333333"

PLOT_COLORS = ["#1F77B4", "#AEC7E8", "#FF7F0E"]
CATEGORY_COLORS = {"DATA":"#1F77B4", "VOICE":"#AEC7E8", "MIX":"#FF7F0E", "SMS":"#8FB3D9", "OTHERS":"#B0B0B0"}

def _plot_colors(labels):
    """Use a sober palette: one main color for packages, limited colors for categories."""
    labels = list(labels or [])
    if labels and all(str(label).upper() in CATEGORY_COLORS for label in labels):
        return [CATEGORY_COLORS[str(label).upper()] for label in labels]
    return [PLOT_COLORS[0] for _ in labels]

def _add_percent_labels_on_bars(ax, bars, values, force_raw_percent=False):
    vals = [float(v or 0) for v in values or []]
    total = sum(abs(v) for v in vals) or 1.0
    y0, y1 = ax.get_ylim()
    pad = (y1 - y0) * 0.02
    for bar, value in zip(bars, vals):
        pct_value = value if force_raw_percent else (value / total * 100)
        label = f"{pct_value:.2f}%"
        x = bar.get_x() + bar.get_width() / 2
        y = bar.get_height()
        ax.text(x, y + (pad if y >= 0 else -pad), label, ha="center", va="bottom" if y >= 0 else "top", fontsize=8, fontweight="normal", color="#333333")


def _slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            parts.append(shape.text)
    return "\n".join(parts)


def _remove_slide_by_index(prs, index: int):
    xml_slides = prs.slides._sldIdLst
    slide_id = list(xml_slides)[index]
    prs.part.drop_rel(slide_id.rId)
    xml_slides.remove(slide_id)


def _keep_only_template_thank_you_slide(prs):
    """Keep the official Chinguitel thank-you slide and remove template cover slides."""
    thank_indexes = []
    for i, slide in enumerate(prs.slides):
        txt = _slide_text(slide).lower()
        if "merci" in txt or "شك" in txt:
            thank_indexes.append(i)
    keep_index = thank_indexes[-1] if thank_indexes else None
    for i in reversed(range(len(prs.slides))):
        if i != keep_index:
            _remove_slide_by_index(prs, i)
    return keep_index is not None


def _move_first_slide_to_end(prs):
    if len(prs.slides) <= 1:
        return
    xml_slides = prs.slides._sldIdLst
    first = list(xml_slides)[0]
    xml_slides.remove(first)
    xml_slides.append(first)


def _add_fallback_thank_you_slide(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb(WHITE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(2.35), Inches(0.03), Inches(0.9))
    line.fill.solid(); line.fill.fore_color.rgb = _rgb(VIOLET); line.line.color.rgb = _rgb(VIOLET)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.45), Inches(5.0), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = "Merci شكرا"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _rgb(BLUE)
    return slide


def _rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))




def _blank_layout(prs):
    """Return a safe blank slide layout for any user template.
    Some PowerPoint templates contain fewer than 7 layouts, so slide_layouts[6]
    can raise IndexError. We prefer a truly blank layout, otherwise use the
    last available layout.
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise IndexError("The PowerPoint template contains no slide layouts.")
    for layout in layouts:
        if str(getattr(layout, "name", "")).strip().lower() == "blank":
            return layout
    return layouts[-1]

def _today_label():
    return timezone.now().strftime("%d/%m/%Y")


def _add_title(slide, title: str, subtitle: str | None = None):
    # Style proche de la template Marketing Weekly Report : titre bleu centré,
    # date en haut à droite, fond blanc et typographie simple.
    date_box = slide.shapes.add_textbox(Inches(9.85), Inches(0.22), Inches(1.45), Inches(0.28))
    dp = date_box.text_frame.paragraphs[0]
    dp.text = _today_label()
    dp.font.size = Pt(9)
    dp.font.color.rgb = _rgb(DARK)
    dp.alignment = PP_ALIGN.RIGHT

    shape = slide.shapes.add_textbox(Inches(0.65), Inches(0.52), Inches(10.4), Inches(0.45))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = _rgb(BLUE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.9), Inches(0.98), Inches(9.9), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.alignment = PP_ALIGN.CENTER
        sp.font.name = "Arial"
        sp.font.size = Pt(10)
        sp.font.color.rgb = _rgb(DARK)


def _add_footer(slide):
    return


def _new_slide(prs, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(_blank_layout(prs))
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb(WHITE)
    _add_title(slide, title, subtitle)
    return slide


def _add_metric(slide, x, y, w, h, label, value):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(WHITE)
    box.line.color.rgb = _rgb(BLUE)
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    p1 = tf.paragraphs[0]
    p1.text = str(value)
    p1.alignment = PP_ALIGN.CENTER
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = _rgb(BLUE)
    p2 = tf.add_paragraph()
    p2.text = label
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(9)
    p2.font.color.rgb = _rgb(DARK)


def _plot_line(labels, values, title="") -> BytesIO:
    fig = plt.figure(figsize=(12.6, 6.15))
    ax = fig.add_subplot(111)
    ax.plot(labels, values, linewidth=2, marker="o", markersize=3, color="#0070C0")
    ax.set_title(title, color="#0070C0", fontweight="bold", fontsize=12)
    ax.tick_params(axis="x", rotation=45, labelsize=8)
    ax.tick_params(axis="y", labelsize=8)
    ax.grid(True, axis="y", alpha=0.25)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    stream = BytesIO()
    fig.savefig(stream, format="png", dpi=260, facecolor="white", bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)
    stream.seek(0)
    return stream




def _plot_multi_line(labels, series, title="") -> BytesIO:
    fig = plt.figure(figsize=(12.6, 6.15))
    ax = fig.add_subplot(111)
    labels = list(labels or [])
    for i, serie in enumerate(series or []):
        name = str(serie.get("name", f"Series {i+1}"))
        values = [float(v or 0) for v in (serie.get("data") or [])]
        ax.plot(labels, values, linewidth=2, marker="o", markersize=3, label=name, color=PLOT_COLORS[i % len(PLOT_COLORS)])
    ax.set_title(title, color="#0070C0", fontweight="bold", fontsize=12)
    ax.tick_params(axis="x", rotation=45, labelsize=7)
    ax.tick_params(axis="y", labelsize=8)
    ax.grid(True, axis="y", alpha=0.25)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    if series:
        ax.legend(fontsize=7, loc="best")
    fig.tight_layout()
    stream = BytesIO()
    fig.savefig(stream, format="png", dpi=260, facecolor="white", bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)
    stream.seek(0)
    return stream

def _plot_bar(labels, values, title="") -> BytesIO:
    fig = plt.figure(figsize=(12.6, 6.15))
    ax = fig.add_subplot(111)
    labels = list(labels or [])
    vals = [float(v or 0) for v in (values or [])]
    bars = ax.bar(labels, vals, color=_plot_colors(labels))
    ax.set_title(title, color="#0070C0", fontweight="bold", fontsize=12)
    ax.tick_params(axis="x", rotation=45, labelsize=8)
    ax.tick_params(axis="y", labelsize=8)
    ax.grid(True, axis="y", alpha=0.25)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    force_raw_percent = "%" in str(title)
    if force_raw_percent:
        ax.yaxis.set_major_formatter(PercentFormatter(xmax=100))
    _add_percent_labels_on_bars(ax, bars, vals, force_raw_percent=force_raw_percent)
    fig.tight_layout()
    stream = BytesIO()
    fig.savefig(stream, format="png", dpi=260, facecolor="white", bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)
    stream.seek(0)
    return stream



def _percent_values(values):
    vals = [float(v or 0) for v in (values or [])]
    total = sum(vals) or 1.0
    return [round(v / total * 100, 2) for v in vals]

def _fmt_pct_value(v):
    if v is None or v == "":
        return "N/A"
    try:
        return f"{float(v):.2f}%"
    except Exception:
        return str(v)

def _rows_share_only(rows):
    converted = []
    for r in rows or []:
        item = dict(r)
        if "share" in item:
            item["share"] = _fmt_pct_value(item.get("share"))
        if "trend" in item:
            item["trend"] = _fmt_pct_value(item.get("trend")) if item.get("trend") is not None else "N/A"
        if "change" in item:
            item["change"] = _fmt_pct_value(item.get("change")) if item.get("change") is not None else "N/A"
        item.pop("revenue", None)
        item.pop("avg", None)
        converted.append(item)
    return converted

def _plot_donut(labels, values, title="") -> BytesIO:
    """Kept name for compatibility, but draws a bar chart instead of a circular chart."""
    fig = plt.figure(figsize=(12.0, 6.15))
    ax = fig.add_subplot(111)
    labels = list(labels or [])
    vals = [float(v or 0) for v in (values or [])]
    bars = ax.bar(labels, vals, color=_plot_colors(labels))
    _add_percent_labels_on_bars(ax, bars, vals)
    ax.set_title(title, color="#0070C0", fontweight="bold", fontsize=12)
    ax.tick_params(axis="x", labelrotation=35, labelsize=8)
    ax.grid(axis="y", alpha=0.25)
    fig.tight_layout()
    stream = BytesIO()
    fig.savefig(stream, format="png", dpi=260, facecolor="white", bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)
    stream.seek(0)
    return stream


def _add_picture(slide, stream, x=0.22, y=1.02, w=12.90, h=6.25):
    """Insert charts almost full-slide for readable PowerPoint export.

    The exported graph keeps only a small top band for the slide title/date and
    uses the rest of the 16:9 slide area. This fixes small charts such as
    "WoW % — Selected Category" when presented in PowerPoint.
    """
    slide.shapes.add_picture(stream, Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def _format_cell_value(value):
    if isinstance(value, float):
        return f"{value:,.2f}"
    if isinstance(value, int):
        return str(value)
    return str(value)

def _add_table(slide, rows, x=0.65, y=1.35, w=10.4, h=5.6, max_rows=10):
    rows = rows[:max_rows]
    if not rows:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        p.text = "Aucune donnée disponible."
        p.font.size = Pt(14)
        return
    headers = list(rows[0].keys())
    # Compact table: calculate a smaller height based on number of rows to avoid huge empty spaces.
    row_count = len(rows) + 1
    compact_h = min(h, max(0.45, 0.34 * row_count))
    table_shape = slide.shapes.add_table(row_count, len(headers), Inches(x), Inches(y), Inches(w), Inches(compact_h))
    table = table_shape.table
    for r in range(row_count):
        table.rows[r].height = Inches(0.34 if r else 0.38)
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(BLUE)
        cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.color.rgb = _rgb(WHITE)
        para.font.bold = True
        para.font.size = Pt(8)
    for i, row in enumerate(rows, start=1):
        for j, header in enumerate(headers):
            cell = table.cell(i, j)
            cell.text = _format_cell_value(row.get(header, ""))
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(8)
            para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(LIGHT)


def _category_comparison_rows(rows, current_key="current", previous_key="previous", current_label="This Week", previous_label="Previous Week"):
    out = []
    total_current = 0.0
    for r in rows or []:
        if str(r.get("category", "")).upper() == "TOTAL":
            continue
        try:
            total_current += float(r.get(current_key) or 0)
        except Exception:
            pass
    for r in rows or []:
        cat = str(r.get("category", ""))
        cur = float(r.get(current_key) or 0)
        prev = float(r.get(previous_key) or 0)
        contrib = (cur / total_current * 100) if total_current and cat.upper() != "TOTAL" else (100.0 if cat.upper() == "TOTAL" and cur else 0.0)
        out.append({"Category": cat, current_label: round(cur, 2), previous_label: round(prev, 2), "Contribution (%)": f"{contrib:.2f}%"})
    return out

def _top_pack_comparison_rows(group, current_label="This Week", previous_label="Previous Week"):
    """Top Packs table with current period, previous period and contribution.

    It uses the union of current and previous packages. This avoids hiding packs
    that are present in the previous week/month but missing or low in the current
    period, such as Beinatna 500/700.
    """
    current = group.get("current_top") or group.get("top_current") or group.get("top_week") or group.get("top_month") or []
    previous = group.get("previous_top") or group.get("previous_month_top") or group.get("previous_week_top") or []
    cur_map = {str(r.get("package")): float(r.get("revenue") or 0) for r in current}
    prev_map = {str(r.get("package")): float(r.get("revenue") or 0) for r in previous}
    all_packages = sorted(set(cur_map) | set(prev_map), key=lambda p: (cur_map.get(p, 0.0), prev_map.get(p, 0.0)), reverse=True)
    current_total = sum(cur_map.values()) or 0.0
    rows = []
    for pkg in all_packages:
        cur = cur_map.get(pkg, 0.0)
        prev = prev_map.get(pkg, 0.0)
        contribution = (cur / current_total * 100) if current_total else 0
        rows.append({
            "Package": pkg,
            current_label: round(cur, 2),
            previous_label: round(prev, 2),
            "Contribution (%)": f"{contribution:.2f}%",
        })
    return rows

def _fmt_num(value, decimals=2):
    try:
        return f"{float(value or 0):,.{decimals}f}"
    except Exception:
        return str(value)

def _add_report_chart_slide(prs, title, labels, values, chart_title=None, kind="bar", subtitle=None):
    slide = _new_slide(prs, title, subtitle)
    if kind == "line":
        _add_picture(slide, _plot_line(labels, values, chart_title or title))
    elif kind == "donut":
        _add_picture(slide, _plot_donut(labels, values, chart_title or title))
    else:
        _add_picture(slide, _plot_bar(labels, values, chart_title or title))
    return slide

def _add_report_multi_line_slide(prs, title, labels, series, subtitle=None):
    slide = _new_slide(prs, title, subtitle)
    _add_picture(slide, _plot_multi_line(labels, series, title))
    return slide

def generate_pptx(filters: dict | None = None, upload: DataUpload | None = None, generated_by_task=False) -> PowerPointReport:
    filters = filters or {}
    if upload and "upload_id" not in filters:
        filters["upload_id"] = str(upload.id)
    qs = filter_records(filters)

    kpi = summary(qs)
    daily = time_series(qs, "day")
    weekly = time_series(qs, "week")
    packages = package_breakdown(qs, 12)
    categories = category_breakdown(qs)
    fcst = forecast_4_weeks(qs)
    anom = detect_anomalies(qs)
    heat = calendar_heatmap(qs)
    eda_data = eda(qs)

    # Optional new analytics are imported lazily to keep old reports compatible.
    from .analytics import demand_dashboard, ml_dashboard, monthly_marketing_report, weekly_marketing_report, daily_marketing_report
    demand_data = demand_dashboard(qs, filters.get("granularity", "week"))
    ml_data = ml_dashboard(qs)
    monthly_data = monthly_marketing_report(qs, filters.get("month_id"))

    # Detail chart selected from the Category detail screen (DATA / VOICE / MIX / ...).
    detail_category = (filters.get("detail_category") or filters.get("category") or "DATA").upper()
    detail_wow_data = None

    # Export dans la template officielle Chinguitel fournie par l'utilisateur.
    # Structure obligatoire du PPT exporté :
    #   1) première slide de la template = couverture
    #   2) toutes les slides générées automatiquement = graphes + tableaux
    #   3) dernière slide de la template = Merci / شكرا
    template_path = Path(settings.BASE_DIR) / "Weekly report Marketing 04-10-May-26 new.pptx"
    thank_slide_element = None
    template_original_count = 0
    template_thank_index = None
    if template_path.exists():
        prs = Presentation(str(template_path))
        template_original_count = len(prs.slides)

        # IMPORTANT FIX:
        # Do NOT delete template slides before adding generated slides.
        # python-pptx can then reuse an existing slide file name (for example slide3.xml),
        # which creates duplicate ZIP entries and PowerPoint opens the file in Repair mode.
        # We keep the template intact while generating slides, then remove the unused
        # original template slides at the end.
        for i, slide0 in enumerate(prs.slides):
            txt = _slide_text(slide0).lower()
            if "merci" in txt or "شك" in txt:
                template_thank_index = i
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    sections = set((filters.get("sections") or "").split(",")) if filters.get("sections") else {"kpi", "daily", "weekly", "monthly", "packages", "categories", "demand", "ml", "forecast", "anomalies", "tables"}
    sections.discard("stats")

    if "detail_wow" in sections:
        detail_filters = dict(filters)
        detail_filters["category"] = detail_category
        detail_wow_data = detail_dashboard(filter_records(detail_filters), "week")

    weekly_report_data = weekly_marketing_report(qs, filters.get("week_id")) if "weekly_report" in sections else None
    daily_report_data = daily_marketing_report(qs, filters.get("week_id")) if "daily_report" in sections else None
    monthly_report_data = monthly_marketing_report(qs, filters.get("month_id")) if "monthly_report" in sections else None

    if "detail_wow" in sections and detail_wow_data:
        dw = detail_wow_data.get("weekly", {})
        slide = _new_slide(prs, "WoW % — Selected Category", f"{detail_category} · Variation vs previous week")
        _add_picture(slide, _plot_bar(dw.get("labels", []), dw.get("wow", []), "WoW % — Selected Category"))

    if "weekly_report" in sections and weekly_report_data:
        wr = weekly_report_data
        k = wr.get("kpis", {})
        slide = _new_slide(prs, "Weekly Report — KPI", wr.get("selected_week", {}).get("label", ""))
        _add_metric(slide, 0.8, 1.5, 3.7, 1.1, "Total Revenue", _fmt_num(k.get("selected_revenue")))
        _add_metric(slide, 4.8, 1.5, 3.7, 1.1, "Last Week", _fmt_num(k.get("previous_revenue")))
        _add_metric(slide, 8.8, 1.5, 3.7, 1.1, "Change Rate", _fmt_pct_value(k.get("change")))
        _add_metric(slide, 0.8, 3.2, 3.7, 1.1, "Average / Week", _fmt_num(k.get("avg_daily")))
        _add_metric(slide, 4.8, 3.2, 3.7, 1.1, "Packages", k.get("packages", 0))

        tr = wr.get("total_revenue", {})
        slide = _new_slide(prs, "Weekly Report — Total Revenue")
        _add_picture(slide, _plot_bar(tr.get("labels", []), tr.get("revenue_millions") or tr.get("revenue", []), "Total Revenue by week"))

        ad = wr.get("avg_daily", {})
        slide = _new_slide(prs, "Weekly Report — Average Revenue per Week")
        _add_picture(slide, _plot_bar(ad.get("labels", []), ad.get("total", []), "Average Revenue per Week"))

        if "tables" in sections:
            rows = wr.get("category_compare", {}).get("rows", [])
            if rows:
                slide = _new_slide(prs, "Weekly Report — Category comparison")
                _add_table(slide, _rows_share_only(rows), max_rows=12)

        for group in (wr.get("groups") or [])[:6]:
            slide = _new_slide(prs, f"Weekly Report — Packs rev — {group.get('name','')}")
            _add_picture(slide, _plot_bar(group.get("labels", []), group.get("revenue_millions") or group.get("revenue", []), f"{group.get('name','')} revenue"))
            if "tables" in sections:
                top_rows = group.get("current_top") or group.get("top_current") or []
                if top_rows:
                    slide = _new_slide(prs, f"Weekly Report — Top packs — {group.get('name','')}")
                    _add_table(slide, _top_pack_comparison_rows(group, "This Week", "Previous Week"), max_rows=12)

    if "daily_report" in sections and daily_report_data:
        dr = daily_report_data
        k = dr.get("kpis", {})
        subtitle = dr.get("selected_week", {}).get("label", "")
        slide = _new_slide(prs, "Daily Report — KPI", subtitle)
        _add_metric(slide, 0.8, 1.5, 3.7, 1.1, "Total Revenue", _fmt_num(k.get("selected_revenue")))
        _add_metric(slide, 4.8, 1.5, 3.7, 1.1, "Average / day", _fmt_num(k.get("avg_daily")))
        _add_metric(slide, 8.8, 1.5, 3.7, 1.1, "Best day", k.get("best_day", "—"))
        _add_metric(slide, 0.8, 3.2, 3.7, 1.1, "Best day revenue", _fmt_num(k.get("best_day_revenue")))
        _add_metric(slide, 4.8, 3.2, 3.7, 1.1, "Mon → Sun change", _fmt_pct_value(k.get("change_monday_sunday")))
        _add_metric(slide, 8.8, 3.2, 3.7, 1.1, "Packages", k.get("packages", 0))

        rv = dr.get("daily_revenue", {})
        _add_report_chart_slide(prs, "Daily Report — Daily revenue", rv.get("labels", []), rv.get("revenue_millions") or rv.get("revenue", []), "Daily Revenue", "bar", subtitle)

        cd = dr.get("category_daily", {})
        _add_report_multi_line_slide(prs, "Daily Report — Category daily trend", cd.get("labels", []), cd.get("series", []), subtitle)
        if "tables" in sections:
            slide = _new_slide(prs, "Daily Report — Category comparison", subtitle)
            _add_table(slide, _category_comparison_rows(dr.get("category_compare", {}).get("rows", []), "total", "previous_total", "This Week", "Previous Week"), max_rows=12)
            slide = _new_slide(prs, "Daily Report — Day KPI table", subtitle)
            _add_table(slide, _rows_share_only(dr.get("day_kpis", [])), max_rows=10)

        for group in (dr.get("groups") or [])[:6]:
            _add_report_chart_slide(prs, f"Daily Report — Packs rev — {group.get('name','')}", group.get("labels", []), group.get("revenue_millions") or group.get("revenue", []), f"{group.get('name','')} daily revenue", "bar", subtitle)
            if "tables" in sections and group.get("top_week"):
                slide = _new_slide(prs, f"Daily Report — Top packs — {group.get('name','')}", subtitle)
                _add_table(slide, _top_pack_comparison_rows(group, "This Week", "Previous Week"), max_rows=12)

        if "tables" in sections:
            for item in (dr.get("day_package_tables") or [])[:7]:
                slide = _new_slide(prs, f"Daily Report — Top packs — {item.get('label','')}", subtitle)
                _add_table(slide, _rows_share_only(item.get("rows", [])), max_rows=12)

    if "monthly_report" in sections and monthly_report_data:
        mrpt = monthly_report_data
        k = mrpt.get("kpis", {})
        subtitle = mrpt.get("selected_month", {}).get("label", "")
        slide = _new_slide(prs, "Monthly Report — KPI", subtitle)
        _add_metric(slide, 0.8, 1.5, 3.7, 1.1, "Total Revenue", _fmt_num(k.get("selected_revenue")))
        _add_metric(slide, 4.8, 1.5, 3.7, 1.1, "Previous Month", _fmt_num(k.get("previous_revenue")))
        _add_metric(slide, 8.8, 1.5, 3.7, 1.1, "Change Rate", _fmt_pct_value(k.get("change")))
        _add_metric(slide, 0.8, 3.2, 3.7, 1.1, "Average / Month", _fmt_num(k.get("avg_daily")))
        _add_metric(slide, 4.8, 3.2, 3.7, 1.1, "Packages", k.get("packages", 0))

        mrev = mrpt.get("monthly_revenue", {})
        _add_report_chart_slide(prs, "Monthly Report — Monthly revenue", mrev.get("labels", []), mrev.get("revenue", []), "Monthly Revenue", "bar", subtitle)
        if "tables" in sections:
            slide = _new_slide(prs, "Monthly Report — Category comparison", subtitle)
            _add_table(slide, _category_comparison_rows(mrpt.get("category_compare", {}).get("rows", []), "current", "previous", "This Month", "Previous Month"), max_rows=12)
        for group in (mrpt.get("groups") or [])[:6]:
            _add_report_chart_slide(prs, f"Monthly Report — Packs rev — {group.get('name','')}", group.get("labels", []), group.get("revenue", []), f"{group.get('name','')} monthly revenue", "bar", subtitle)
            if "tables" in sections and group.get("top_month"):
                slide = _new_slide(prs, f"Monthly Report — Top packs — {group.get('name','')}", subtitle)
                _add_table(slide, _top_pack_comparison_rows(group, "This Month", "Previous Month"), max_rows=12)

    if "kpi" in sections:
        slide = _new_slide(prs, "KPIs principaux")
        _add_metric(slide, 0.8, 1.5, 3.7, 1.1, "Revenu total", _fmt_num(kpi["total_revenue"]))
        _add_metric(slide, 4.8, 1.5, 3.7, 1.1, "Moyenne quotidienne", _fmt_num(kpi["daily_average"]))
        _add_metric(slide, 8.8, 1.5, 3.7, 1.1, "Packages distincts", kpi["packages"])
        _add_metric(slide, 0.8, 3.2, 3.7, 1.1, "Nombre de lignes", kpi["rows"])
        _add_metric(slide, 4.8, 3.2, 3.7, 1.1, "Nombre de jours", kpi["days"])
        _add_metric(slide, 8.8, 3.2, 3.7, 1.1, "WoW %", f"{kpi['wow_percent']}%" if kpi["wow_percent"] is not None else "N/A")

    if "demand" in sections:
        slide = _new_slide(prs, "Dashboard principal — demande par forfait (%)")
        d = demand_data.get("demand_ranking", {})
        _add_picture(slide, _plot_bar(d.get("labels", []), d.get("share", []), "Ranking des forfaits les plus populaires (%)"))
        slide = _new_slide(prs, "Pareto 80/20 — revenus")
        pr = demand_data.get("pareto", {})
        _add_picture(slide, _plot_bar(pr.get("labels", []), _percent_values(pr.get("revenue", [])), "Distribution des revenus par forfait (%)"))
        pop = demand_data.get("popularity_evolution", {})
        if pop.get("series"):
            _add_report_multi_line_slide(prs, "Popularity evolution — top packs (%)", pop.get("labels", []), pop.get("series", []))
        can = demand_data.get("cannibalization", {})
        if can.get("series"):
            _add_report_multi_line_slide(prs, "Cannibalisation — package comparison", can.get("labels", []), can.get("series", []))
        avg = demand_data.get("revenue_avg", {})
        _add_report_chart_slide(prs, "Average revenue by package", avg.get("labels", []), avg.get("avg", []), "Average Revenue by Package", "bar")
        per = demand_data.get("period_revenue", {})
        _add_report_chart_slide(prs, "Revenue by period (%)", per.get("labels", []), per.get("share", []), "Revenue by period (%)", "bar")
        if "tables" in sections:
            slide = _new_slide(prs, "Segmentation des forfaits par performance")
            _add_table(slide, _rows_share_only(demand_data.get("performance_segments", [])), max_rows=12)
            slide = _new_slide(prs, "Cash cows / Stars")
            _add_table(slide, _rows_share_only(demand_data.get("cash_cows", [])), max_rows=12)
            slide = _new_slide(prs, "Low profitability products")
            _add_table(slide, _rows_share_only(demand_data.get("low_profit", [])), max_rows=12)

    if "daily" in sections:
        slide = _new_slide(prs, "Évolution journalière")
        _add_picture(slide, _plot_line(daily["labels"], daily["revenue"], "Revenu par jour"))

    if "weekly" in sections:
        slide = _new_slide(prs, "Vue hebdomadaire")
        _add_picture(slide, _plot_bar(weekly["labels"], weekly["revenue"], "Revenu par semaine"))

    if "monthly" in sections:
        slide = _new_slide(prs, "Monthly Report")
        mr = monthly_data.get("monthly_revenue", {})
        _add_picture(slide, _plot_bar(mr.get("labels", []), mr.get("revenue", []), "Revenu mensuel"))
        if "tables" in sections:
            slide = _new_slide(prs, "Monthly Report — comparaison par catégorie")
            _add_table(slide, _rows_share_only(monthly_data.get("category_compare", {}).get("rows", [])), max_rows=10)

    if "packages" in sections:
        slide = _new_slide(prs, "Top packages")
        _add_picture(slide, _plot_bar(packages["labels"], _percent_values(packages["revenue"]), "Top packages — part du revenu (%)"))

    if "categories" in sections:
        slide = _new_slide(prs, "Répartition par catégorie")
        _add_picture(slide, _plot_bar(categories["labels"], _percent_values(categories["revenue"]), "Part du revenu par catégorie (%)"))

    if "ml" in sections:
        if "anomalies" in sections:
            slide = _new_slide(prs, "Machine Learning — anomalies")
            _add_table(slide, ml_data.get("anomalies", {}).get("rows", []), max_rows=12)
        slide = _new_slide(prs, "Machine Learning — clustering")
        _add_table(slide, ml_data.get("clusters", []), max_rows=8)
        slide = _new_slide(prs, "Machine Learning — trend analysis")
        _add_table(slide, ml_data.get("trends", []), max_rows=12)

    if "forecast" in sections:
        slide = _new_slide(prs, "Prévisions 3 semaines", f"Méthode : {ml_data.get('forecast', {}).get('method') or fcst.get('method')}")
        f_points = ml_data.get("forecast", {}).get("points", []) or fcst.get("points", [])[:21]
        f_labels = [p["date"] for p in f_points]
        f_values = [p["yhat"] for p in f_points]
        _add_picture(slide, _plot_line(f_labels, f_values, "Forecast journalier — 21 jours"))

    if "anomalies" in sections and "ml" not in sections:
        slide = _new_slide(prs, "Anomalies détectées")
        _add_table(slide, anom["rows"], max_rows=12)

    # Heatmap et EDA supprimés du PowerPoint

    # Nettoyer la template seulement maintenant (après création des slides générées).
    # On garde la couverture (slide 0), la slide Merci originale si elle existe,
    # et toutes les slides générées. Cette méthode évite les fichiers ppt/slides/*.xml
    # dupliqués qui déclenchent l'erreur "PowerPoint a détecté un problème".
    if template_path.exists() and template_original_count:
        generated_indexes = set(range(template_original_count, len(prs.slides)))
        keep = {0} | generated_indexes
        if template_thank_index is not None:
            keep.add(template_thank_index)
        for i in reversed(range(template_original_count)):
            if i not in keep:
                _remove_slide_by_index(prs, i)

        # Remettre la slide Merci originale exactement à la fin.
        thank_slide_element = None
        for i, slide0 in enumerate(prs.slides):
            txt = _slide_text(slide0).lower()
            if "merci" in txt or "شك" in txt:
                thank_slide_element = list(prs.slides._sldIdLst)[i]
        if thank_slide_element is not None:
            try:
                prs.slides._sldIdLst.remove(thank_slide_element)
                prs.slides._sldIdLst.append(thank_slide_element)
            except ValueError:
                pass
        else:
            _add_fallback_thank_you_slide(prs)

    reports_dir = Path(settings.MEDIA_ROOT) / "generated_reports"
    reports_dir.mkdir(parents=True, exist_ok=True)
    with NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        tmp_path = Path(tmp.name)

    report = PowerPointReport(upload=upload, filters=filters, generated_by_task=generated_by_task)
    filename = f"chinguitel_revenue_report_{timezone.now():%Y%m%d_%H%M%S}.pptx"
    with tmp_path.open("rb") as f:
        report.file.save(filename, File(f), save=True)
    tmp_path.unlink(missing_ok=True)
    return report
